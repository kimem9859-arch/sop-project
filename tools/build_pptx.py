# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
CREAM   = RGBColor(0xFD,0xF4,0xED)
CREAMD  = RGBColor(0xFB,0xE8,0xD8)
ORANGE  = RGBColor(0xE8,0x5D,0x2F)
ORANGEL = RGBColor(0xF4,0xA2,0x61)
RED     = RGBColor(0xD6,0x30,0x27)
REDD    = RGBColor(0xB8,0x24,0x1C)
INK     = RGBColor(0x2B,0x2B,0x2B)
GRAY    = RGBColor(0x6B,0x6B,0x6B)
GRAYL   = RGBColor(0xD8,0xD2,0xCC)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
FONT    = "Malgun Gothic"

SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def set_ea(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin','a:ea','a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', name)

def tb(slide, x, y, w, h, paras, anchor='t', wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {'t':MSO_ANCHOR.TOP,'m':MSO_ANCHOR.MIDDLE,'b':MSO_ANCHOR.BOTTOM}[anchor]
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        if 'sa' in p: para.space_after = Pt(p['sa'])
        if 'sb' in p: para.space_before = Pt(p['sb'])
        if 'ls' in p: para.line_spacing = p['ls']
        for r in p['runs']:
            run = para.add_run(); run.text = r['text']
            f = run.font
            f.size = Pt(r.get('size',14)); f.bold = r.get('bold',False); f.italic = r.get('italic',False)
            f.color.rgb = r.get('color', INK)
            set_ea(run, r.get('font', FONT))
    return box

def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, rounded=False, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if shadow:
        sp = shp._element.spPr
        ef = sp.makeelement(qn('a:effectLst'), {})
        sh = sp.makeelement(qn('a:outerShdw'), {'blurRad':'90000','dist':'30000','dir':'5400000','rotWithShape':'0'})
        clr = sp.makeelement(qn('a:srgbClr'), {'val':'B8241C'})
        al = sp.makeelement(qn('a:alpha'), {'val':'10000'})
        clr.append(al); sh.append(clr); ef.append(sh); sp.append(ef)
    return shp

def oval(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp

def bg(slide, color=CREAM):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def deco(slide):
    # faint warm circle top-right
    oval(slide, SW-2.6, -1.7, 4.3, 4.3, CREAMD)

def header(slide, num, ko, en):
    tx = 1.45
    if num:
        rect(slide, 0.62, 0.5, 0.62, 0.62, fill=ORANGE, rounded=True)
        tb(slide, 0.62, 0.5, 0.62, 0.62, [{'runs':[{'text':num,'size':22,'bold':True,'color':WHITE}],'align':PP_ALIGN.CENTER}], anchor='m')
    else:
        tx = 0.62
    tb(slide, tx, 0.46, 10.5, 0.6, [{'runs':[{'text':ko,'size':27,'bold':True,'color':INK}]}], anchor='m')
    tb(slide, tx+0.02, 1.08, 10.5, 0.32, [{'runs':[{'text':en,'size':11,'bold':True,'color':ORANGE}]}])

def footer(slide, page):
    tb(slide, 0.62, 7.02, 6.0, 0.3, [{'runs':[{'text':'Vision AI 웨어러블 안전 콘솔  ·  4조','size':9,'color':GRAYL}]}])
    tb(slide, SW-2.0, 7.0, 1.4, 0.3, [{'runs':[{'text':page,'size':10,'bold':True,'color':GRAY}],'align':PP_ALIGN.RIGHT}])

def content_slide(num, ko, en, page):
    s = prs.slides.add_slide(BLANK); bg(s); deco(s); header(s, num, ko, en); footer(s, page)
    return s

# =================== S1 표지 ===================
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, 0.28, SH, fill=ORANGE)                      # left accent band
rect(s, 0.28, 0, 0.12, SH, fill=RED)
oval(s, SW-3.2, -2.2, 5.6, 5.6, CREAMD)
# tags
tb(s, 1.1, 1.05, 11, 0.4, [{'runs':[
    {'text':'#VisionAI   ','size':12,'bold':True,'color':ORANGE},
    {'text':'#EdgeAI   ','size':12,'bold':True,'color':ORANGE},
    {'text':'#Wearable   ','size':12,'bold':True,'color':ORANGE},
    {'text':'#FSM   ','size':12,'bold':True,'color':ORANGE},
    {'text':'#Fail-Safe','size':12,'bold':True,'color':ORANGE}]}])
# big title
tb(s, 1.1, 1.75, 11, 2.0, [
    {'runs':[{'text':'작업자 Human Error','size':46,'bold':True,'color':INK}],'sa':4},
    {'runs':[{'text':'사전 예방 ','size':46,'bold':True,'color':INK},
             {'text':'안전 콘솔','size':46,'bold':True,'color':RED}]}])
rect(s, 1.13, 3.78, 2.4, 0.06, fill=ORANGE)
tb(s, 1.1, 4.0, 11, 0.5, [{'runs':[{'text':'Vision AI 기반 웨어러블 연동 Fail-Safe 안전 시스템','size':17,'bold':True,'color':GRAY}]}])
tb(s, 1.1, 4.5, 11, 0.4, [{'runs':[{'text':'단일 PECVD 콘솔의 공정 순서 위반(누락·역순)을 버튼 입력 전 사전 차단하는 PoC','size':12,'color':GRAY}]}])
# team block
rect(s, 1.1, 5.35, 11.1, 1.35, fill=WHITE, line=GRAYL, lw=0.75, rounded=True)
oval(s, 1.4, 5.62, 0.8, 0.8, ORANGE)
tb(s, 1.4, 5.62, 0.8, 0.8, [{'runs':[{'text':'4조','size':16,'bold':True,'color':WHITE}],'align':PP_ALIGN.CENTER}], anchor='m')
tb(s, 2.45, 5.55, 6.5, 1.0, [
    {'runs':[{'text':'한국폴리텍대학 청주캠퍼스 · 반도체시스템과 · 2학년 B반','size':12,'bold':True,'color':INK}],'sa':3},
    {'runs':[{'text':'지도교수  허주회','size':11,'color':GRAY}],'sa':3},
    {'runs':[{'text':'김응민 · 김동현 · 박송빈 · 이재모 · 김태현 · 신희재 · 김선원 · 천희동','size':10,'color':GRAY}]}], anchor='m')
tb(s, 9.0, 5.55, 3.0, 1.0, [
    {'runs':[{'text':'2026. 06. 24','size':15,'bold':True,'color':RED}],'align':PP_ALIGN.RIGHT,'sa':3},
    {'runs':[{'text':'1학기 기말 발표','size':11,'color':GRAY}],'align':PP_ALIGN.RIGHT}], anchor='m')
tb(s, 1.1, 6.95, 11, 0.3, [{'runs':[{'text':'2026학년도 융합프로젝트실습  |  한국폴리텍대학 청주캠퍼스','size':9,'color':GRAYL}]}])

# =================== S2 목차 ===================
s = content_slide('', '목차', 'CONTENTS', '02')
# override header number area: 목차 uses no number chip nicely; keep orange chip empty-> replace with dot
items = [
    ('01','프로젝트 개요','PROJECT OVERVIEW'),
    ('02','배경 및 필요성','BACKGROUND & NEED'),
    ('03','AS-IS / TO-BE','REACTION → PREVENTION'),
    ('04','요구사항 정의','REQUIREMENTS'),
    ('05','시스템 구성 · 시연 시나리오','SYSTEM & DEMO'),
    ('06','핵심 기술 (AI · FSM)','CORE TECHNOLOGY'),
    ('07','성능 검증 · 진행 현황','PERFORMANCE & STATUS'),
    ('08','향후 계획','ROADMAP'),
]
colx = [0.9, 6.95]
for i,(n,ko,en) in enumerate(items):
    cx = colx[i//4]; ry = 1.95 + (i%4)*1.18
    rect(s, cx, ry, 5.5, 1.0, fill=WHITE, line=GRAYL, lw=0.75, rounded=True)
    tb(s, cx+0.22, ry, 0.9, 1.0, [{'runs':[{'text':n,'size':24,'bold':True,'color':ORANGEL}]}], anchor='m')
    tb(s, cx+1.15, ry, 4.2, 1.0, [
        {'runs':[{'text':ko,'size':15,'bold':True,'color':INK}],'sa':2},
        {'runs':[{'text':en,'size':9,'bold':True,'color':GRAY}]}], anchor='m')

# =================== S3 프로젝트 개요 ===================
s = content_slide('01', '프로젝트 개요', 'PROJECT OVERVIEW', '03')
# left: one-line definition card
rect(s, 0.62, 1.75, 6.0, 4.7, fill=ORANGE, rounded=True, shadow=True)
tb(s, 1.0, 2.1, 5.3, 0.5, [{'runs':[{'text':'한 줄 정의','size':13,'bold':True,'color':RGBColor(0xFF,0xE3,0xD5)}]}])
tb(s, 1.0, 2.65, 5.3, 3.4, [
    {'runs':[{'text':'단일 PECVD 콘솔','size':22,'bold':True,'color':WHITE}],'sa':6,'ls':1.1},
    {'runs':[{'text':'에서 공정 기동 시퀀스의 ','size':16,'color':WHITE},
             {'text':'작업 순서 위반(누락·역순)','size':16,'bold':True,'color':RGBColor(0xFF,0xE3,0xD5)},
             {'text':'을 비전 AI로 감지해,','size':16,'color':WHITE}],'sa':6,'ls':1.25},
    {'runs':[{'text':'오조작을 ','size':16,'color':WHITE},
             {'text':'버튼 입력 전에 사전 차단','size':16,'bold':True,'color':WHITE},
             {'text':'하는 능동형 Fail-Safe 안전 시스템','size':16,'color':WHITE}],'ls':1.25}], anchor='t')
tb(s, 1.0, 5.75, 5.3, 0.6, [{'runs':[
    {'text':'⚠ 만능 아님 — 단일 공정 순서 위반 1종 ','size':11,'bold':True,'color':RGBColor(0xFF,0xE3,0xD5)},
    {'text':'PoC','size':11,'bold':True,'color':WHITE}]}], anchor='m')
# right: key-value facts
facts = [
    ('소속','한국폴리텍대 청주캠 · 반도체시스템과 · 4조 (8인)'),
    ('지도교수','허주회'),
    ('트랙','융합프로젝트(학교 캡스톤, 본체) + 한이음(예산·공모전 연계)'),
    ('기간','2026.02 ~ 2026.10.30'),
    ('핵심 키워드','#VisionAI #EdgeAI #Wearable #FSM #Fail-Safe'),
]
ry = 1.75
for k,v in facts:
    rect(s, 6.95, ry, 5.75, 0.86, fill=WHITE, line=GRAYL, lw=0.75, rounded=True)
    rect(s, 6.95, ry, 0.09, 0.86, fill=ORANGE)
    tb(s, 7.25, ry, 1.7, 0.86, [{'runs':[{'text':k,'size':11,'bold':True,'color':ORANGE}]}], anchor='m')
    tb(s, 8.95, ry, 3.6, 0.86, [{'runs':[{'text':v,'size':11,'color':INK}],'ls':1.05}], anchor='m')
    ry += 0.96

# =================== S4 배경 및 필요성 ===================
s = content_slide('02', '배경 및 필요성', 'BACKGROUND & NEED', '04')
cards = [
    ('1','휴먼 에러는 산업의 큰 문제','제조 품질 결함의 약 80%가 휴먼 에러에서 비롯 (NIST·IJERA 기준). 다룰 가치가 충분한 문제.'),
    ('2','기존 장비의 사후 반응 한계','버튼이 눌린 "사후"에만 반응 → 오작동한 뒤에야 인지. 사람 주의력에만 의존하는 구조의 한계.'),
    ('3','작업 순서 위반의 위험','누락·역순은 앞 단계의 안전 전제가 무너진 채 다음 동작 실행 → 장비 오작동·공정 불량·설비 손상.'),
]
cw = 3.95; gap = 0.2; x0 = 0.62
for i,(n,t,d) in enumerate(cards):
    cx = x0 + i*(cw+gap)
    rect(s, cx, 1.9, cw, 3.05, fill=WHITE, line=GRAYL, lw=0.75, rounded=True, shadow=True)
    oval(s, cx+0.35, 2.25, 0.72, 0.72, ORANGE)
    tb(s, cx+0.35, 2.25, 0.72, 0.72, [{'runs':[{'text':n,'size':22,'bold':True,'color':WHITE}],'align':PP_ALIGN.CENTER}], anchor='m')
    tb(s, cx+0.35, 3.2, cw-0.7, 0.85, [{'runs':[{'text':t,'size':15,'bold':True,'color':INK}],'ls':1.1}])
    tb(s, cx+0.35, 4.0, cw-0.7, 0.85, [{'runs':[{'text':d,'size':11,'color':GRAY}],'ls':1.25}])
# bottom callout
rect(s, 0.62, 5.2, 12.1, 1.25, fill=CREAMD, rounded=True)
rect(s, 0.62, 5.2, 0.12, 1.25, fill=RED)
tb(s, 1.0, 5.4, 11.5, 0.9, [
    {'runs':[{'text':'핵심 메시지   ','size':13,'bold':True,'color':RED},
             {'text':'순서 검증을 사람의 집중력에만 맡길 수 없다.','size':14,'bold':True,'color':INK}],'sa':4},
    {'runs':[{'text':'→ 사후 경고 위에, 버튼 입력 직전 자동으로 감지·차단하는 사전 예방형 안전 장치(인터록)가 필요하다.','size':12,'color':GRAY}]}], anchor='m')

# =================== S5 AS-IS / TO-BE ===================
s = content_slide('03', 'AS-IS / TO-BE', 'FROM REACTION TO PREVENTION', '05')
rows = [('반응 시점','버튼이 눌린 "사후"에만 반응\n(오작동 후 인지)','버튼 입력 "직전"에 사전 감지·차단\n레이어를 추가'),
        ('순서 검증','순서 위반을 작업자의\n주의력에 의존','비전 AI가 ROI 기준으로\n자동 판정 (정답/오답)'),
        ('사고 대응','오작동 발생 후\n보고·조사 (사후 처리)','인터록으로 전기 신호를\n원천 차단 (사전 차단)')]
# AS-IS card (left, gray) / TO-BE card (right, orange)
rect(s, 0.62, 1.85, 5.9, 3.5, fill=WHITE, line=GRAYL, lw=1.0, rounded=True, shadow=True)
rect(s, 6.83, 1.85, 5.9, 3.5, fill=WHITE, line=ORANGE, lw=1.5, rounded=True, shadow=True)
rect(s, 0.62, 1.85, 5.9, 0.62, fill=GRAY, rounded=True)
rect(s, 0.62, 2.16, 5.9, 0.31, fill=GRAY)
rect(s, 6.83, 1.85, 5.9, 0.62, fill=ORANGE, rounded=True)
rect(s, 6.83, 2.16, 5.9, 0.31, fill=ORANGE)
tb(s, 0.62, 1.85, 5.9, 0.62, [{'runs':[{'text':'AS-IS  ·  기존 콘솔','size':15,'bold':True,'color':WHITE}],'align':PP_ALIGN.CENTER}], anchor='m')
tb(s, 6.83, 1.85, 5.9, 0.62, [{'runs':[{'text':'TO-BE  ·  본 안전 콘솔','size':15,'bold':True,'color':WHITE}],'align':PP_ALIGN.CENTER}], anchor='m')
ry = 2.65
for k,a,b in rows:
    tb(s, 0.95, ry, 1.35, 0.8, [{'runs':[{'text':k,'size':11,'bold':True,'color':GRAY}]}], anchor='m')
    tb(s, 2.35, ry, 4.0, 0.8, [{'runs':[{'text':a.replace('\n',' '),'size':12,'color':INK}],'ls':1.1}], anchor='m')
    tb(s, 7.15, ry, 1.35, 0.8, [{'runs':[{'text':k,'size':11,'bold':True,'color':ORANGE}]}], anchor='m')
    tb(s, 8.55, ry, 4.0, 0.8, [{'runs':[{'text':b.replace('\n',' '),'size':12,'bold':True,'color':INK}],'ls':1.1}], anchor='m')
    if ry < 4.4:
        rect(s, 0.95, ry+0.83, 5.25, 0.012, fill=GRAYL)
        rect(s, 7.15, ry+0.83, 5.25, 0.012, fill=CREAMD)
    ry += 0.9
# callout +alpha
rect(s, 0.62, 5.62, 12.1, 1.05, fill=CREAMD, rounded=True)
rect(s, 0.62, 5.62, 0.12, 1.05, fill=ORANGE)
tb(s, 1.0, 5.62, 11.5, 1.05, [{'runs':[
    {'text':'핵심 메시지  ＋α   ','size':13,'bold':True,'color':RED},
    {'text':'기존 콘솔을 대체·부정하지 않는다. 그 "사후 반응" 위에 ','size':13,'color':INK},
    {'text':'버튼 입력 직전의 사전 감지·차단 레이어를 얹는','size':13,'bold':True,'color':INK},
    {'text':' 개선이다.','size':13,'color':INK}],'ls':1.2}], anchor='m')

prs.save('/home/kimem/projects/기말발표_5장.pptx')
print('saved 기말발표_5장.pptx  slides:', len(prs.slides._sldIdLst))
